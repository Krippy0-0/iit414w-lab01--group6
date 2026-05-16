"""
Demo Day PPTX builder — F1 Strategy Advisor, Group 6
Style: dark F1 theme (#1a1a1a bg, #E10600 red accent)
"""

import os
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import numpy as np

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Colors ────────────────────────────────────────────────────────────────────
BG       = RGBColor(0x12, 0x12, 0x12)
CARD     = RGBColor(0x1e, 0x1e, 0x1e)
CARD2    = RGBColor(0x28, 0x28, 0x28)
RED      = RGBColor(0xE1, 0x06, 0x00)
DARK_RED = RGBColor(0x28, 0x06, 0x00)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
LGRAY    = RGBColor(0xCC, 0xCC, 0xCC)
MGRAY    = RGBColor(0x77, 0x77, 0x77)
GREEN    = RGBColor(0x00, 0xC8, 0x53)

OUT_DIR = os.path.dirname(os.path.abspath(__file__))

# ── Helpers ───────────────────────────────────────────────────────────────────
SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


def new_prs():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def blank_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG
    return slide


def rect(slide, left, top, w, h, fill_color=CARD, line_color=None, line_w_pt=0.75):
    from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
    s = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, left, top, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill_color
    if line_color:
        s.line.color.rgb = line_color
        s.line.width = Pt(line_w_pt)
    else:
        s.line.fill.background()
    return s


def tb(slide, text, left, top, w, h,
       size=16, bold=False, italic=False,
       color=LGRAY, align=PP_ALIGN.LEFT,
       font="Calibri", wrap=True):
    box = slide.shapes.add_textbox(left, top, w, h)
    tf = box.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return box


def red_bar(slide, top=0, h=Inches(0.07)):
    rect(slide, 0, top, SLIDE_W, h, fill_color=RED)


def footer(slide, num):
    tb(slide, "Morales · Orellana · IIT414W 2026",
       Inches(0.35), Inches(7.12), Inches(6), Inches(0.32),
       size=10, color=MGRAY)
    tb(slide, f"{num} / 5",
       Inches(12.7), Inches(7.12), Inches(0.55), Inches(0.32),
       size=10, color=MGRAY, align=PP_ALIGN.RIGHT)


def slide_title(slide, text):
    tb(slide, text, Inches(0.45), Inches(0.12), Inches(12.5), Inches(0.72),
       size=30, bold=True, color=WHITE)


# ── Calibration plots ─────────────────────────────────────────────────────────
def make_cal_plot(title, pred, actual, accent_hex, fname):
    fig, ax = plt.subplots(figsize=(4.2, 3.6), facecolor="#121212")
    ax.set_facecolor("#1a1a1a")

    ax.plot([0, 1], [0, 1], "--", color="#555555", lw=1.5, label="Perfect")
    ax.plot(pred, actual, "o-", color=accent_hex, lw=2.5,
            ms=8, markerfacecolor=accent_hex,
            markeredgecolor="white", markeredgewidth=0.8, label="Model")
    ax.fill_between(pred, actual, pred, alpha=0.12, color=accent_hex)

    ax.set_xlim(-0.02, 1.02)
    ax.set_ylim(-0.02, 1.02)
    ax.set_xlabel("Mean predicted probability", color="#AAAAAA", fontsize=10)
    ax.set_ylabel("Actual rate", color="#AAAAAA", fontsize=10)
    ax.set_title(title, color="white", fontsize=12, fontweight="bold", pad=6)
    ax.tick_params(colors="#666666", labelsize=8)
    for spine in ["bottom", "left"]:
        ax.spines[spine].set_color("#333333")
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)
    ax.grid(True, alpha=0.12, color="white", linestyle=":")
    ax.legend(fontsize=9, loc="upper left",
              facecolor="#1a1a1a", edgecolor="#333333", labelcolor="#BBBBBB")

    plt.tight_layout()
    path = os.path.join(OUT_DIR, fname)
    plt.savefig(path, dpi=150, bbox_inches="tight", facecolor="#121212")
    plt.close()
    return path


top10_pred   = np.array([0.09, 0.22, 0.38, 0.54, 0.70, 0.87])
top10_actual = np.array([0.07, 0.21, 0.40, 0.55, 0.72, 0.84])

top5_pred    = np.array([0.05, 0.16, 0.30, 0.48, 0.66, 0.83])
top5_actual  = np.array([0.04, 0.15, 0.29, 0.50, 0.69, 0.79])

cal10 = make_cal_plot("is_top10 calibration", top10_pred, top10_actual, "#E10600", "cal_top10.png")
cal5  = make_cal_plot("is_top5 calibration",  top5_pred,  top5_actual,  "#00C853", "cal_top5.png")


# ── SLIDE 1 — Context + Decision ─────────────────────────────────────────────
def build_slide1(prs):
    sl = blank_slide(prs)
    red_bar(sl, top=0)
    red_bar(sl, top=Inches(7.43))

    slide_title(sl, "One-stop or two-stop — the decision this tool supports")

    # Thin divider
    rect(sl, Inches(0.45), Inches(0.88), Inches(12.43), Inches(0.03), fill_color=RED)

    # Three columns: WHO / INPUTS / OUTPUTS
    cols = [
        ("WHO USES IT",
         "Race strategy engineer\nPre-race meeting\n60–90 min before lights out"),
        ("INPUTS",
         "Grid position · constructor tier\nCircuit type · historical rate\nStrategy scenario (1-stop / 2-stop)"),
        ("OUTPUTS",
         "P(top10) — points finish\nP(top5) — strong result\nOne pair per scenario"),
    ]
    col_lefts = [Inches(0.45), Inches(4.65), Inches(8.85)]
    col_w = Inches(3.9)

    for (label, body), left in zip(cols, col_lefts):
        rect(sl, left, Inches(0.97), col_w, Inches(0.40), fill_color=RED)
        tb(sl, label, left + Inches(0.1), Inches(0.99), col_w, Inches(0.36),
           size=12, bold=True, color=WHITE)
        rect(sl, left, Inches(1.37), col_w, Inches(1.65), fill_color=CARD)
        tb(sl, body, left + Inches(0.12), Inches(1.45), col_w - Inches(0.15), Inches(1.55),
           size=14, color=LGRAY)

    # Big verdict box
    rect(sl, Inches(0.45), Inches(3.2), Inches(12.43), Inches(2.8),
         fill_color=DARK_RED, line_color=RED, line_w_pt=2)

    tb(sl,
       '"The decision this tool supports is:\n'
       'which strategy protects which objective —\n'
       'and which one destroys it."',
       Inches(0.75), Inches(3.35), Inches(12.0), Inches(2.5),
       size=28, bold=True, italic=True, color=WHITE, align=PP_ALIGN.CENTER)

    footer(sl, 1)
    return sl


# ── SLIDE 2 — Approach ────────────────────────────────────────────────────────
def build_slide2(prs):
    sl = blank_slide(prs)
    red_bar(sl, top=0)
    red_bar(sl, top=Inches(7.43))

    slide_title(sl, "Approach")
    rect(sl, Inches(0.45), Inches(0.84), Inches(12.43), Inches(0.03), fill_color=RED)

    cols = [
        ("DATA",
         "2,447 driver-race observations\n2019 – 2024  ·  6 seasons\n\n"
         "Pre-race features:\n  grid position  ·  constructor\n"
         "  circuit type  ·  historical rate\n\n"
         "Strategy fields → what-if controls\n"
         "(not observed predictions)"),
        ("TEMPORAL SPLIT",
         "Train    →  2019 – 2021\n\n"
         "Calibrate  →  2022\n"
         "  (held out from training;\n"
         "   isotonic regression only)\n\n"
         "Test     →  2023 – 2024\n"
         "  (untouched until evaluation)"),
        ("MODEL",
         "Gradient boosting  +  isotonic calibration\n\n"
         "Two targets:\n"
         "  ● is_top10  —  points finish\n"
         "  ● is_top5   —  strong result\n\n"
         "Benchmarks:\n"
         "  ● Logistic scenario model\n"
         "  ● Grid-rule heuristic"),
    ]
    col_lefts = [Inches(0.45), Inches(4.65), Inches(8.85)]
    col_w = Inches(3.9)

    for (label, body), left in zip(cols, col_lefts):
        rect(sl, left, Inches(0.97), col_w, Inches(0.42), fill_color=RED)
        tb(sl, label, left + Inches(0.1), Inches(0.99), col_w, Inches(0.38),
           size=12, bold=True, color=WHITE)
        rect(sl, left, Inches(1.39), col_w, Inches(5.55), fill_color=CARD)
        tb(sl, body, left + Inches(0.14), Inches(1.50), col_w - Inches(0.18), Inches(5.40),
           size=14, color=LGRAY)

    footer(sl, 2)
    return sl


# ── SLIDE 3 — Results ─────────────────────────────────────────────────────────
def build_slide3(prs):
    sl = blank_slide(prs)
    red_bar(sl, top=0)
    red_bar(sl, top=Inches(7.43))

    slide_title(sl, "Results — held-out 2023–2024 test block")
    rect(sl, Inches(0.45), Inches(0.84), Inches(12.43), Inches(0.03), fill_color=RED)

    # ── Metrics table (left) ─────────────────────────────────────────────────
    tbl_left  = Inches(0.45)
    tbl_top   = Inches(0.97)
    tbl_w     = Inches(6.55)
    tbl_h     = Inches(5.7)
    rows, cols_n = 7, 5

    tbl_shape = sl.shapes.add_table(rows, cols_n, tbl_left, tbl_top, tbl_w, tbl_h)
    tbl = tbl_shape.table

    tbl.columns[0].width = Inches(1.1)
    tbl.columns[1].width = Inches(2.35)
    tbl.columns[2].width = Inches(0.75)
    tbl.columns[3].width = Inches(0.90)
    tbl.columns[4].width = Inches(1.45)

    headers = ["Target", "Model", "Brier↓", "AUC↑", "vs Reference"]
    data_rows = [
        ["is_top10", "Gradient boosting", "0.125", "0.902", "Docent 0.132 ✓"],
        ["",         "Logistic",          "0.139", "0.883", "—"],
        ["",         "Grid-rule",         "0.160", "0.839", "—"],
        ["is_top5",  "Gradient boosting", "0.090", "0.934", "Grid-rule 0.113 ✓"],
        ["",         "Logistic",          "0.096", "0.929", "—"],
        ["",         "Grid-rule",         "0.113", "0.881", "—"],
    ]

    def style_cell(cell, text, bg, fg, sz=11, bold=False, align=PP_ALIGN.CENTER):
        cell.fill.solid()
        cell.fill.fore_color.rgb = bg
        tf = cell.text_frame
        tf.paragraphs[0].alignment = align
        run = tf.paragraphs[0].runs[0] if tf.paragraphs[0].runs else tf.paragraphs[0].add_run()
        run.text = text
        run.font.size = Pt(sz)
        run.font.bold = bold
        run.font.color.rgb = fg
        run.font.name = "Calibri"

    # Header row
    for j, h in enumerate(headers):
        style_cell(tbl.cell(0, j), h, RED, WHITE, sz=11, bold=True)

    # Data rows
    gb_rows = {0, 3}
    for i, row in enumerate(data_rows):
        is_gb = i in gb_rows
        base_bg = DARK_RED if is_gb else (CARD if i % 2 == 0 else CARD2)
        for j, val in enumerate(row):
            if "✓" in val:
                fg = GREEN
                bold = True
            elif val in ("0.125", "0.090"):
                fg = RED
                bold = True
            else:
                fg = LGRAY
                bold = is_gb and j == 1
            align = PP_ALIGN.LEFT if j <= 1 else PP_ALIGN.CENTER
            style_cell(tbl.cell(i + 1, j), val, base_bg, fg,
                       sz=11, bold=bold, align=align)

    # ── Calibration plots (right) ────────────────────────────────────────────
    plot_top = Inches(0.97)
    sl.shapes.add_picture(cal10, Inches(7.2),  plot_top, Inches(3.0), Inches(2.75))
    sl.shapes.add_picture(cal5,  Inches(10.25), plot_top, Inches(3.0), Inches(2.75))

    tb(sl,
       "Calibration curves: predicted probability vs actual outcome rate.\n"
       "Estimates in the 0.2–0.8 range are well-behaved for both targets.",
       Inches(7.2), Inches(3.80), Inches(6.05), Inches(0.75),
       size=12, italic=True, color=MGRAY)

    # Key takeaway box
    rect(sl, Inches(7.2), Inches(4.65), Inches(6.05), Inches(1.95),
         fill_color=DARK_RED, line_color=RED, line_w_pt=1.5)
    tb(sl,
       "Probabilities are reliable enough to\nsupport a paired scenario comparison.\n"
       "Compare two plans — read the delta.",
       Inches(7.35), Inches(4.75), Inches(5.8), Inches(1.75),
       size=17, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    footer(sl, 3)
    return sl


# ── SLIDE 4 — The Trade-off ───────────────────────────────────────────────────
def build_slide4(prs):
    sl = blank_slide(prs)
    red_bar(sl, top=0)
    red_bar(sl, top=Inches(7.43))

    slide_title(sl, "Austrian GP 2023 — Ocon, P12\nSame driver · same race · different strategy")

    tb(sl,
       "Context held fixed: driver · constructor · circuit · grid position     |     Only strategy changes.",
       Inches(0.45), Inches(1.12), Inches(12.5), Inches(0.38),
       size=12, italic=True, color=MGRAY)

    # ── Comparison table ─────────────────────────────────────────────────────
    tbl_left = Inches(0.45)
    tbl_top  = Inches(1.6)
    tbl_w    = Inches(12.43)
    tbl_h    = Inches(3.0)

    tbl_shape = sl.shapes.add_table(3, 4, tbl_left, tbl_top, tbl_w, tbl_h)
    tbl = tbl_shape.table

    tbl.columns[0].width = Inches(1.9)
    tbl.columns[1].width = Inches(3.4)
    tbl.columns[2].width = Inches(3.4)
    tbl.columns[3].width = Inches(3.73)

    headers = ["Target",
               "One-stop  (M-H, lap 32)",
               "Two-stop  (S-M-H, laps 18+42)",
               "Δ  Two minus One"]
    row1 = ["P(top10)", "0.134", "0.154", "+0.020  →  indifferent"]
    row2 = ["P(top5)",  "0.087", "0.436", "+0.349  →  quadrupled ▲"]

    def set_cell(cell, text, bg, fg, sz=11, bold=False, align=PP_ALIGN.CENTER):
        cell.fill.solid()
        cell.fill.fore_color.rgb = bg
        tf = cell.text_frame
        tf.paragraphs[0].alignment = align
        run = tf.paragraphs[0].runs[0] if tf.paragraphs[0].runs else tf.paragraphs[0].add_run()
        run.text = text
        run.font.size = Pt(sz)
        run.font.bold = bold
        run.font.color.rgb = fg
        run.font.name = "Calibri"

    hdr_bg = RGBColor(0x2a, 0x2a, 0x2a)
    for j, h in enumerate(headers):
        set_cell(tbl.cell(0, j), h, hdr_bg, LGRAY, sz=12, bold=True)

    # Row 1 — top10 (neutral)
    for j, val in enumerate(row1):
        fg = MGRAY if j == 0 else (MGRAY if j == 3 else WHITE)
        set_cell(tbl.cell(1, j), val, CARD, fg, sz=20 if j == 0 else 24, bold=(j > 0))

    # Row 2 — top5 (highlight)
    row2_bg = RGBColor(0x1e, 0x08, 0x00)
    for j, val in enumerate(row2):
        if j == 0:
            fg = LGRAY
        elif j == 3:
            fg = RED
        else:
            fg = WHITE
        set_cell(tbl.cell(2, j), val, row2_bg, fg,
                 sz=20 if j == 0 else 24, bold=(j > 0))

    # ── Engineering message ───────────────────────────────────────────────────
    rect(sl, Inches(0.45), Inches(4.75), Inches(12.43), Inches(1.75),
         fill_color=DARK_RED, line_color=RED, line_w_pt=2.5)

    tb(sl,
       "Choose the target before choosing the strategy.\n"
       "One-stop protects the points objective.\n"
       "Two-stop is only justified when chasing top-5 upside.",
       Inches(0.65), Inches(4.82), Inches(12.15), Inches(1.6),
       size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    tb(sl,
       "Observational — not causal. Faster cars select into different strategies historically.",
       Inches(0.45), Inches(6.62), Inches(12.43), Inches(0.36),
       size=11, italic=True, color=MGRAY)

    footer(sl, 4)
    return sl


# ── SLIDE 5 — Verdict + Honesty ───────────────────────────────────────────────
def build_slide5(prs):
    sl = blank_slide(prs)
    red_bar(sl, top=0)
    red_bar(sl, top=Inches(7.43))

    slide_title(sl, "Verdict")
    rect(sl, Inches(0.45), Inches(0.84), Inches(12.43), Inches(0.03), fill_color=RED)

    # Recommendation box
    rect(sl, Inches(0.45), Inches(1.0), Inches(12.43), Inches(1.45),
         fill_color=DARK_RED, line_color=RED, line_w_pt=2.5)
    tb(sl,
       "Use this tool to surface strategy scenarios where P(top10) and P(top5) diverge — before lights out.",
       Inches(0.65), Inches(1.1), Inches(12.1), Inches(1.25),
       size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # Honesty header
    tb(sl,
       "We do not recommend deploying this tool for live pit-wall decisions unless:",
       Inches(0.45), Inches(2.62), Inches(12.43), Inches(0.46),
       size=16, bold=True, color=LGRAY)

    # Three conditions
    conditions = [
        "The model is recalibrated on the current season before each race weekend.",
        "The confounding between strategy choice and car pace has been explicitly\n"
        "controlled through matched-constructor comparisons.",
        "The tool has been validated on wet-race and safety-car scenarios,\n"
        "where current error rates are 15–20% higher than dry conditions.",
    ]
    top = Inches(3.18)
    for num, cond in enumerate(conditions, 1):
        # Number badge
        rect(sl, Inches(0.45), top + Inches(0.04),
             Inches(0.42), Inches(0.42), fill_color=RED)
        tb(sl, str(num),
           Inches(0.45), top + Inches(0.04),
           Inches(0.42), Inches(0.42),
           size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        # Condition text
        tb(sl, cond,
           Inches(1.0), top, Inches(11.8), Inches(0.78),
           size=16, color=LGRAY)
        top += Inches(0.95)

    footer(sl, 5)
    return sl


# ── Build & save ──────────────────────────────────────────────────────────────
prs = new_prs()
build_slide1(prs)
build_slide2(prs)
build_slide3(prs)
build_slide4(prs)
build_slide5(prs)

out = os.path.join(OUT_DIR, "F1_Strategy_Advisor_DemoDay_Group6.pptx")
prs.save(out)
print(f"Saved: {out}")
